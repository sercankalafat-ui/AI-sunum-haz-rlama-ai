from __future__ import annotations

import tempfile
from pathlib import Path

from .analyzer import AnalysisResult


def _build_score_chart(result: AnalysisResult) -> Path:
    try:
        import matplotlib.pyplot as plt
    except ModuleNotFoundError as error:
        raise RuntimeError("Grafik için matplotlib kurulmalı: pip install matplotlib") from error

    products = [str(item["product"]) for item in result.ranked_products]
    scores = [float(item["overall_score"]) for item in result.ranked_products]

    figure, axis = plt.subplots(figsize=(10, 5))
    axis.bar(products, scores, color="#2E86AB")
    axis.set_title("Ürün Genel Skor Karşılaştırması")
    axis.set_ylabel("Skor")
    axis.set_ylim(0, 1)
    axis.grid(axis="y", alpha=0.3)

    output = Path(tempfile.gettempdir()) / "product_score_chart.png"
    figure.tight_layout()
    figure.savefig(output, dpi=200)
    plt.close(figure)
    return output


def create_presentation(
    result: AnalysisResult,
    output_path: str = "urun-analizi-sunumu.pptx",
) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ModuleNotFoundError as error:
        raise RuntimeError("PowerPoint üretimi için python-pptx kurulmalı: pip install python-pptx") from error

    presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Ürün Analizi ve Öneri Sunumu"
    title_slide.placeholders[1].text = "Yapay zeka destekli otomatik hazırlık"

    summary_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    summary_slide.shapes.title.text = "Genel Durum Özeti"
    summary_box = summary_slide.placeholders[1].text_frame
    summary_box.text = f"Ortalama fiyat: {result.summary['average_price']:.2f}"
    for line in [
        f"Ortalama kalite: {result.summary['average_quality']:.2f}",
        f"Ortalama müşteri memnuniyeti: {result.summary['average_customer_satisfaction']:.2f}",
        f"Ortalama satış: {result.summary['average_sales']:.2f}",
    ]:
        summary_box.add_paragraph().text = line

    ranking_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    ranking_slide.shapes.title.text = "Skor Sıralaması"
    table = ranking_slide.shapes.add_table(
        rows=len(result.ranked_products) + 1,
        cols=3,
        left=Inches(0.5),
        top=Inches(1.4),
        width=Inches(9),
        height=Inches(4.5),
    ).table
    table.cell(0, 0).text = "Ürün"
    table.cell(0, 1).text = "Genel Skor"
    table.cell(0, 2).text = "Satış"

    for index, item in enumerate(result.ranked_products, start=1):
        table.cell(index, 0).text = str(item["product"])
        table.cell(index, 1).text = f"{item['overall_score']:.3f}"
        table.cell(index, 2).text = str(item["sales"])

    chart_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    chart_slide.shapes.title.text = "Skor Grafiği"
    chart_image = _build_score_chart(result)
    chart_slide.shapes.add_picture(str(chart_image), Inches(0.7), Inches(1.3), width=Inches(8.8))

    recommendation_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    recommendation_slide.shapes.title.text = "Öneriler"
    top = result.ranked_products[0]
    recommendation_slide.placeholders[1].text = (
        f"En yüksek potansiyele sahip ürün: {top['product']} (Skor: {top['overall_score']:.3f}).\n"
        "Satış-kalite-memnuniyet dengesini koruyarak bu ürün odağında kampanya yürütülmesi önerilir."
    )

    out = Path(output_path)
    presentation.save(out)
    return out
